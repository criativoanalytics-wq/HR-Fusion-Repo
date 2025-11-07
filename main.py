from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from io import BytesIO
import tempfile
import docx2txt
from PyPDF2 import PdfReader
import spacy
from datetime import datetime
import tempfile as tmp
import os, re, json
import numpy as np
from sentence_transformers import SentenceTransformer
import faiss
import docx

# ============================================================
# 🚀 AIDA DRIVE CONNECTOR - RAG VERSION (Multilíngue e Smart)
# ============================================================

app = FastAPI(
    title="AIDA Drive Connector",
    description="API RAG multilíngue para leitura e busca semântica no Google Drive (.docx, .pdf, .txt)",
    version="2.1.0"
)

# ============================================================
# 🌐 CORS
# ============================================================
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
# Carrega modelos multilíngues
nlp_en = spacy.load("en_core_web_sm")
nlp_pt = spacy.load("pt_core_news_sm")

# ============================================================
# 🔐 Autenticação
# ============================================================
def get_service():
    """Cria o serviço autenticado do Google Drive."""
    if not os.path.exists("token.json"):
        raise HTTPException(status_code=401, detail="Token OAuth ausente. Gere o token primeiro com auth_setup.py")
    creds = Credentials.from_authorized_user_file("token.json", SCOPES)
    return build("drive", "v3", credentials=creds)

# ============================================================
# 🧠 Dicionário de sinônimos bilíngue
# ============================================================
SINONIMOS = {
    "governança de dados": ["data governance", "gestão de dados", "política de dados", "data management"],
    "qualidade de dados": ["data quality", "data cleansing", "data validation"],
    "catálogo de dados": ["data catalog", "metadata management"],
    "lago de dados": ["data lake", "data repository"],
    "segurança da informação": ["information security", "data privacy", "cybersecurity"],
    "arquitetura de dados": ["data architecture", "data modeling", "data structure"],
    "integração de dados": ["data integration", "ETL", "data ingestion"],
    "governança": ["governance", "management", "oversight"],
}
# Detecta nomes próprios comuns nos documentos
#NOMES_PESSOAS = ["lisa", "rick", "felipe", "sanders", "gavin", "jennifer"]
def detectar_pessoa_spacy(texto: str):
    """Detecta automaticamente nomes de pessoas em PT/EN usando spaCy."""
    if not texto:
        return []
    pessoas = set()

    # Análise em ambos os idiomas
    for nlp in [nlp_en, nlp_pt]:
        doc = nlp(texto)
        for ent in doc.ents:
            if ent.label_ == "PERSON":
                pessoas.add(ent.text.strip())

    return list(pessoas)

def expandir_termos(query: str):
    """Expande automaticamente termos equivalentes em PT/EN e gera busca case-insensitive."""
    if not query:
        return []

    query_lower = query.lower().strip()
    termos_expandidos = {query_lower}

    for chave, sinonimos in SINONIMOS.items():
        if chave in query_lower or any(s in query_lower for s in sinonimos):
            termos_expandidos.add(chave)
            termos_expandidos.update(sinonimos)

    # Garante unicidade
    return list(set(termos_expandidos))

# ============================================================
# 📁 Listagem de arquivos (com expansão bilíngue)
# ============================================================
# ============================================================
# 📁 Listagem de arquivos (com expansão bilíngue + paginação)
# ============================================================
@app.get("/files")
def listar_arquivos(pasta_id: str = None, query: str = None):
    """
    Lista todos os arquivos do Google Drive (com paginação).
    - Expande automaticamente termos bilíngues.
    - Percorre todas as páginas (sem limite de 100 arquivos).
    - Retorna lista completa com metadados.
    """
    try:
        service = get_service()
        termos_busca = expandir_termos(query)
        if not termos_busca:
            termos_busca = [query.lower()] if query else []

        arquivos_encontrados = []
        ids_vistos = set()

        for termo in termos_busca or [""]:
            q = []
            if pasta_id:
                q.append(f"'{pasta_id}' in parents")
            if termo:
                q.append(f"name contains '{termo}'")
            q.append("trashed=false")
            query_final = " and ".join(q)

            page_token = None
            while True:
                results = service.files().list(
                    q=query_final,
                    fields="nextPageToken, files(id, name, mimeType, modifiedTime, parents)",
                    pageSize=100,
                    pageToken=page_token
                ).execute()

                for f in results.get("files", []):
                    if f["id"] not in ids_vistos:
                        arquivos_encontrados.append(f)
                        ids_vistos.add(f["id"])

                page_token = results.get("nextPageToken")
                if not page_token:
                    break  # ✅ todas as páginas lidas

        return {"arquivos": arquivos_encontrados, "total": len(arquivos_encontrados)}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao listar arquivos: {e}")


@app.get("/smart_read")
def smart_read(file_id: str, query: str):
    """
    Busca leve dentro de um PPTX do Google Drive.
    Retorna apenas os slides que contêm o termo especificado.
    Ideal para agentes com limite de payload.
    """
    try:
        from pptx import Presentation
        from io import BytesIO
        import re, tempfile, os

        if not query:
            raise HTTPException(status_code=400, detail="Parâmetro 'query' é obrigatório.")

        service = get_service()
        file = service.files().get(fileId=file_id, fields="name, mimeType").execute()
        nome = file["name"]
        mime = file["mimeType"]

        if mime not in [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint"
        ]:
            raise HTTPException(status_code=400, detail="O arquivo não é um PowerPoint (.pptx).")

        # 📦 Download temporário
        request = service.files().get_media(fileId=file_id)
        fh = BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
        fh.seek(0)

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(fh.read())
            tmp_path = tmp.name

        prs = Presentation(tmp_path)
        os.remove(tmp_path)

        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            textos = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    textos.append(shape.text.strip())
            full_text = " ".join(textos)
            full_text = re.sub(r"\s{2,}", " ", full_text).strip()

            titulo = textos[0] if textos else f"Slide {i}"
            slides.append({
                "slide_numero": i,
                "titulo": titulo,
                "conteudo": full_text
            })

        # 🔍 Busca simples (case-insensitive)
        query_regex = re.compile(re.escape(query), re.IGNORECASE)
        resultados = [
            s for s in slides
            if query_regex.search(s["conteudo"])
        ]

        if not resultados:
            return {
                "arquivo": nome,
                "query": query,
                "total_slides": len(slides),
                "resultados": [],
                "mensagem": "Nenhum slide contém o termo buscado."
            }

        # 🔎 Resumo dos resultados (leve para GPT)
        return {
            "arquivo": nome,
            "query": query,
            "total_slides": len(slides),
            "slides_encontrados": len(resultados),
            "resultados": resultados[:10]  # limite de segurança
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao executar smart_read: {e}")

@app.get("/search_transcripts")
def search_transcripts(query: str, top_k: int = 5):
    """
    Busca semântica dentro dos transcripts previamente indexados (FAISS + Sentence Transformers).
    Retorna os trechos mais relevantes encontrados no índice local.
    """
    import os, json
    import numpy as np
    import faiss
    from sentence_transformers import SentenceTransformer

    try:
        index_path = "index_cache/transcripts.index"
        meta_path = "index_cache/transcripts_meta.json"

        # 🔹 Verifica se o índice existe
        if not os.path.exists(index_path) or not os.path.exists(meta_path):
            raise HTTPException(
                status_code=404,
                detail="O índice FAISS ainda não foi criado. Execute /index_transcripts primeiro."
            )

        # 🔹 Carrega modelo e índice
        model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2")
        index = faiss.read_index(index_path)

        with open(meta_path, "r", encoding="utf-8") as f:
            metadados = json.load(f)

        # 🔹 Gera embedding da query
        query_embedding = model.encode([query])
        query_embedding = np.array(query_embedding, dtype=np.float32)

        # 🔹 Realiza busca
        distances, indices = index.search(query_embedding, top_k)
        resultados = []
        for idx, dist in zip(indices[0], distances[0]):
            if idx < len(metadados):
                item = metadados[idx]
                resultados.append({
                    "arquivo": item["arquivo"],
                    "file_id": item["file_id"],
                    "trecho": item["trecho"],
                    "similaridade": float(1 - dist / 2)  # escala ~0-1
                })

        if not resultados:
            return {"mensagem": "Nenhum trecho relevante encontrado para essa busca."}

        return {
            "query": query,
            "total_resultados": len(resultados),
            "resultados": resultados
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao buscar no índice de transcripts: {e}")


@app.get("/smart_search")
def smart_search(query: str):
    """
    Realiza uma busca expandida:
    - Foco inicial em arquivos relacionados a uma pessoa específica (se aplicável)
    - Expansão semântica multilíngue caso nenhum resultado direto seja encontrado
    """
    try:
        service = get_service()
        termos = expandir_termos(query)
        pessoas = detectar_pessoa_spacy(query)
        pessoa = pessoas[0].lower() if pessoas else None
        arquivos_final = []
        ids_vistos = set()

        def buscar(termos_busca, foco_pessoa=False):
            encontrados = []
            for termo in termos_busca:
                q = f"name contains '{termo}' and trashed=false"
                results = service.files().list(
                    q=q,
                    fields="files(id, name, mimeType, modifiedTime)",
                    pageSize=100
                ).execute()
                for f in results.get("files", []):
                    if f["id"] not in ids_vistos:
                        try:
                            conteudo = ler_arquivo(f["id"])["conteudo"].lower()
                            # 🔍 filtro por pessoa, se houver
                            if foco_pessoa and pessoa not in f["name"].lower() and pessoa not in conteudo:
                                continue
                            # 🔍 busca textual
                            if any(t in conteudo for t in termos_busca):
                                encontrados.append(f)
                                ids_vistos.add(f["id"])
                        except Exception as err:
                            print(f"Erro ao ler {f['name']}: {err}")
                            continue
            return encontrados

        # 🔹 Nível 1: busca restrita à pessoa
        if pessoa:
            arquivos_final = buscar(termos, foco_pessoa=True)

        # 🔹 Nível 2: expansão se nada for encontrado
        expanded = False
        if not arquivos_final:
            arquivos_final = buscar(termos, foco_pessoa=False)
            expanded = True if pessoa else False

        return {
            "query_original": query,
            "pessoa_detectada": pessoa,
            "busca_expandida": expanded,
            "termos_expandidos": termos,
            "arquivos_encontrados": arquivos_final,
            "total": len(arquivos_final)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro na busca expandida: {e}")

# ============================================================
# 📄 Leitura e extração de conteúdo
# ============================================================
# ============================================================
# 📄 Leitura e extração de conteúdo (DOCX, PDF, TXT, PPTX)
# ============================================================
@app.get("/files/{file_id}")
def ler_arquivo(file_id: str, range_inicio: int = 1, range_fim: int = 15):
    """
    Faz download e extrai texto de arquivos (.docx, .pdf, .txt, .pptx),
    com suporte a leitura paginada, chunking para textos longos e tolerância a falhas.
    Ideal para leitura de transcripts e apresentações extensas.
    """
    import tempfile as tmp
    import re, os
    from pptx import Presentation

    try:
        service = get_service()
        file = service.files().get(fileId=file_id, fields="name, mimeType").execute()
        nome = file["name"]
        mime = file["mimeType"]

        request = service.files().get_media(fileId=file_id)
        fh = BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
        fh.seek(0)
        texto_extraido = ""

        # --------------------------------------------------------
        # 🧩 DOCX (com chunking inteligente)
        # --------------------------------------------------------
        if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            import docx

            try:
                with tmp.NamedTemporaryFile(delete=False, suffix=".docx") as temp_file:
                    temp_file.write(fh.read())
                    temp_path = temp_file.name

                doc = docx.Document(temp_path)

                # 🔹 Extrai todos os parágrafos legíveis
                paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                if not paragraphs:
                    texto_extraido = "⚠️ O arquivo DOCX foi encontrado, mas não contém texto legível."
                else:
                    # 🔹 Junta parágrafos e aplica chunking
                    chunk_size = 5000  # caracteres por bloco
                    chunks = []
                    buffer = ""
                    for paragraph in paragraphs:
                        if len(buffer) + len(paragraph) < chunk_size:
                            buffer += paragraph + "\n"
                        else:
                            chunks.append(buffer.strip())
                            buffer = paragraph + "\n"
                    if buffer:
                        chunks.append(buffer.strip())

                    texto_extraido = "\n".join(chunks[:10])  # limite leve (~50k caracteres)

                    return {
                        "nome": nome,
                        "tipo": mime,
                        "total_chunks": len(chunks),
                        "tamanho_medio": int(sum(len(c) for c in chunks) / max(1, len(chunks))),
                        "conteudo": chunks[:10]  # primeiros 10 blocos
                    }

            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)

        # --------------------------------------------------------
        # 📘 PDF
        # --------------------------------------------------------
        elif mime == "application/pdf":
            reader = PdfReader(fh)
            texto_extraido = "\n".join([p.extract_text() or "" for p in reader.pages])

        # --------------------------------------------------------
        # 📄 TXT
        # --------------------------------------------------------
        elif "text" in mime:
            texto_extraido = fh.read().decode("utf-8", errors="ignore")

        # --------------------------------------------------------
        # 🖼️ PPTX (PowerPoint) — leitura hierárquica com paginação
        # --------------------------------------------------------
        elif mime in [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint"
        ]:
            try:
                with tmp.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
                    temp_file.write(fh.read())
                    temp_path = temp_file.name

                prs = Presentation(temp_path)
                total_slides = len(prs.slides)
                range_inicio = max(1, min(range_inicio, total_slides))
                range_fim = max(range_inicio, min(range_fim, total_slides))

                slides_estruturados = []

                def extrair_texto_recursivo(shape, elementos):
                    if hasattr(shape, "text") and shape.text.strip():
                        elementos.append({
                            "texto": shape.text.strip(),
                            "x": int(getattr(shape, "left", 0)),
                            "y": int(getattr(shape, "top", 0))
                        })
                    if hasattr(shape, "shapes"):
                        for subshape in shape.shapes:
                            extrair_texto_recursivo(subshape, elementos)
                    if hasattr(shape, "has_table") and shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells if cell.text.strip()
                            )
                            if row_text:
                                elementos.append({
                                    "texto": row_text,
                                    "x": int(getattr(shape, "left", 0)),
                                    "y": int(getattr(shape, "top", 0))
                                })

                for i in range(range_inicio, range_fim + 1):
                    slide = prs.slides[i - 1]
                    elementos = []
                    for shape in slide.shapes:
                        extrair_texto_recursivo(shape, elementos)

                    elementos_ordenados = sorted(elementos, key=lambda e: (e["y"], e["x"]))
                    linha_id, ultima_y = 0, None
                    for el in elementos_ordenados:
                        if ultima_y is None or abs(el["y"] - ultima_y) > 200000:
                            linha_id += 1
                            ultima_y = el["y"]
                        el["linha_visual"] = linha_id

                    faixas = {}
                    for e in elementos_ordenados:
                        faixa_id = e["linha_visual"]
                        faixas.setdefault(faixa_id, []).append(e["texto"])
                    faixas_agrupadas = [
                        {"linha_visual": k, "conteudo": " ".join(v)} for k, v in faixas.items()
                    ]

                    notas = ""
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notas = slide.notes_slide.notes_text_frame.text.strip()

                    titulo_slide = next(
                        (e["texto"] for e in elementos_ordenados if e["linha_visual"] == 1),
                        f"Slide {i}"
                    )

                    slides_estruturados.append({
                        "slide_numero": i,
                        "titulo": titulo_slide,
                        "faixas": faixas_agrupadas,
                        "notas": notas
                    })

                texto_extraido = ""
                for s in slides_estruturados:
                    texto_extraido += f"\n\n=== SLIDE {s['slide_numero']} - {s['titulo']} ===\n"
                    for f in s["faixas"]:
                        texto_extraido += f"[Faixa {f['linha_visual']}] {f['conteudo']}\n"

                texto_extraido = (
                    texto_extraido.replace("\\n", "\n")
                    .replace("\r", "\n")
                    .replace("\t", " ")
                )
                texto_extraido = re.sub(r"[\u0000-\u001F\u007F-\u009F]", " ", texto_extraido)
                texto_extraido = re.sub(r"\u00A0", " ", texto_extraido)
                texto_extraido = re.sub(r"\s{2,}", " ", texto_extraido)
                texto_extraido = re.sub(r"\n{2,}", "\n", texto_extraido).strip()

            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)

            return {
                "nome": nome,
                "tipo": mime,
                "intervalo_lido": f"{range_inicio}-{range_fim}",
                "total_slides": total_slides,
                "conteudo": texto_extraido[:80000],
                "conteudo_estruturado": slides_estruturados
            }

        # --------------------------------------------------------
        # ❗ Outros formatos
        # --------------------------------------------------------
        else:
            texto_extraido = f"O tipo de arquivo {mime} não é suportado para leitura direta."

        if not isinstance(texto_extraido, str) or not texto_extraido.strip():
            texto_extraido = "⚠️ O arquivo foi encontrado, mas parece não conter texto legível."

        return {
            "nome": nome,
            "tipo": mime,
            "conteudo": texto_extraido[:150000]
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao ler arquivo: {e}")

@app.get("/index_transcripts")
def indexar_transcripts(pasta_raiz: str = None):
    """
    Indexa automaticamente todos os transcripts (.docx e Google Docs) do Google Drive.
    - Inclui arquivos com nomes como 'transcript', 'meeting', 'reunião', 'minutes', 'notes', etc.
    - Caso nenhum arquivo com essas palavras seja encontrado, indexa todos os .docx da pasta.
    - Cria índice FAISS com embeddings (Sentence Transformers) para busca semântica posterior.
    """

    import time
    inicio = time.time()

    try:
        service = get_service()
        model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2")

        print("🚀 Iniciando indexação de transcripts...")

        termos_reuniao = ["transcript", "meeting", "reunião", "minutes", "call", "discussion", "notes"]
        q_filter = " or ".join([f"name contains '{t}'" for t in termos_reuniao])
        q = f"({q_filter}) and trashed=false"

        if pasta_raiz:
            q += f" and '{pasta_raiz}' in parents"

        # 🔍 Busca inicial (termos de reunião)
        results = service.files().list(
            q=q,
            fields="files(id, name, mimeType, modifiedTime)",
            pageSize=1000
        ).execute()

        arquivos = results.get("files", [])
        if not arquivos:
            print("⚠️ Nenhum transcript com palavras-chave encontrado. Tentando indexar todos os .docx da pasta...")
            q = f"mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document' and trashed=false"
            if pasta_raiz:
                q += f" and '{pasta_raiz}' in parents"

            results = service.files().list(
                q=q,
                fields="files(id, name, mimeType, modifiedTime)",
                pageSize=1000
            ).execute()
            arquivos = results.get("files", [])

        if not arquivos:
            print("❌ Nenhum arquivo .docx encontrado para indexar.")
            return {"status": "⚠️ Nenhum transcript .docx ou Google Docs encontrado."}

        os.makedirs("index_cache", exist_ok=True)
        embeddings_list, metadados = [], []

        total = len(arquivos)
        print(f"📁 {total} arquivos candidatos encontrados.")

        for idx, f in enumerate(arquivos, 1):
            nome = f["name"]
            mime = f["mimeType"]
            file_id = f["id"]
            print(f"\n[{idx}/{total}] 📄 Processando: {nome}")

            try:
                # 🔸 Download ou exportação
                if mime == "application/vnd.google-apps.document":
                    request = service.files().export_media(
                        fileId=file_id,
                        mimeType="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )
                elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    request = service.files().get_media(fileId=file_id)
                else:
                    print(f"⏭️ Tipo ignorado: {mime}")
                    continue

                fh = BytesIO()
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while not done:
                    status, done = downloader.next_chunk()
                fh.seek(0)

                with tmp.NamedTemporaryFile(delete=False, suffix=".docx") as temp_file:
                    temp_file.write(fh.read())
                    temp_path = temp_file.name

                # 🔹 Leitura do documento
                doc = docx.Document(temp_path)
                os.remove(temp_path)

                paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                if not paragraphs:
                    print("⚠️ Documento vazio, ignorado.")
                    continue

                # 🔹 Chunking (divisão por blocos de texto)
                chunk_size = 5000
                buffer = ""
                chunks_arquivo = 0

                for paragraph in paragraphs:
                    if len(buffer) + len(paragraph) < chunk_size:
                        buffer += paragraph + "\n"
                    else:
                        emb = model.encode(buffer)
                        embeddings_list.append(emb)
                        metadados.append({
                            "arquivo": nome,
                            "file_id": file_id,
                            "trecho": buffer[:300] + "..."
                        })
                        chunks_arquivo += 1
                        buffer = paragraph + "\n"

                if buffer:
                    emb = model.encode(buffer)
                    embeddings_list.append(emb)
                    metadados.append({
                        "arquivo": nome,
                        "file_id": file_id,
                        "trecho": buffer[:300] + "..."
                    })
                    chunks_arquivo += 1

                print(f"✅ {chunks_arquivo} trechos indexados de {nome}")

            except Exception as e:
                print(f"⚠️ Erro ao processar {nome}: {e}")
                continue

        if not embeddings_list:
            print("❌ Nenhum trecho válido encontrado para indexação.")
            return {"status": "Nenhum trecho válido encontrado para indexação."}

        # 🔹 Criação do índice FAISS
        embeddings = np.array(embeddings_list, dtype=np.float32)
        index = faiss.IndexFlatL2(embeddings.shape[1])
        index.add(embeddings)

        index_path = "index_cache/transcripts.index"
        meta_path = "index_cache/transcripts_meta.json"

        faiss.write_index(index, index_path)
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(metadados, f, ensure_ascii=False, indent=2)

        tempo = round(time.time() - inicio, 2)
        print(f"\n🏁 Indexação concluída em {tempo}s.")
        print(f"📊 Arquivos: {len(arquivos)} | Trechos: {len(embeddings_list)}")

        return {
            "status": "✅ Indexação concluída com sucesso.",
            "total_arquivos": len(arquivos),
            "total_chunks": len(embeddings_list),
            "index_path": index_path,
            "meta_path": meta_path,
            "tempo_execucao_seg": tempo
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao indexar transcripts: {e}")


@app.get("/index_drive")
def indexar_drive(pasta_raiz: str = None):
    """
    🔄 Indexa todo o conteúdo do Google Drive (recursivamente).
    - Lista todas as pastas e arquivos com paginação.
    - Armazena metadados em drive_index.json para uso rápido.
    """
    try:
        service = get_service()
        arquivos_indexados = []
        pastas_a_visitar = [pasta_raiz] if pasta_raiz else [ "root" ]

        def listar_conteudo(pasta_id, caminho_atual=""):
            page_token = None
            while True:
                results = service.files().list(
                    q=f"'{pasta_id}' in parents and trashed=false",
                    fields="nextPageToken, files(id, name, mimeType, modifiedTime, parents)",
                    pageSize=100,
                    pageToken=page_token
                ).execute()

                for item in results.get("files", []):
                    caminho = f"{caminho_atual}/{item['name']}".strip("/")
                    item["path"] = caminho
                    arquivos_indexados.append(item)

                    # Se for pasta → adiciona à fila
                    if item["mimeType"] == "application/vnd.google-apps.folder":
                        listar_conteudo(item["id"], caminho)

                page_token = results.get("nextPageToken")
                if not page_token:
                    break

        # 🔁 Inicia varredura
        for pasta_id in pastas_a_visitar:
            listar_conteudo(pasta_id)

        # 📦 Cria diretório local para índice
        os.makedirs("index_cache", exist_ok=True)
        index_path = f"index_cache/drive_index.json"

        with open(index_path, "w", encoding="utf-8") as f:
            json.dump({
                "timestamp": datetime.utcnow().isoformat(),
                "total_itens": len(arquivos_indexados),
                "arquivos": arquivos_indexados
            }, f, ensure_ascii=False, indent=2)

        return {
            "status": "✅ Indexação concluída com sucesso",
            "total_arquivos": len(arquivos_indexados),
            "index_path": index_path
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao indexar o Drive: {e}")
# ============================================================
# 🔍 Endpoint raiz
# ============================================================
@app.get("/")
def root():
    return {"message": "✅ AIDA Drive Connector RAG (multilíngue) está ativo e pronto para uso."}

if __name__ == "__main__":
    import uvicorn
    import os

    port = int(os.getenv("PORT", 8080))  # 👈 Render injeta a variável PORT
    uvicorn.run("main:app", host="0.0.0.0", port=port, reload=False)
